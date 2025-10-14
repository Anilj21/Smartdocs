import os
from typing import List, Tuple, Dict
import json
import subprocess
import httpx

import fitz  # PyMuPDF
import pdfplumber
from docx import Document
from pptx import Presentation

from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter


def extract_text(filepath: str) -> str:
	text = ""
	file_ext = (filepath.split(".")[-1] or "").lower()
	if file_ext == "pdf":
		try:
			with pdfplumber.open(filepath) as pdf:
				for page in pdf.pages:
					page_text = page.extract_text() or ""
					text += page_text + "\n"
		except Exception:
			with fitz.open(filepath) as doc:
				for page in doc:
					text += page.get_text() + "\n"
	elif file_ext == "docx":
		doc = Document(filepath)
		for para in doc.paragraphs:
			text += para.text + "\n"
	elif file_ext == "pptx":
		prs = Presentation(filepath)
		for slide in prs.slides:
			for shape in slide.shapes:
				if hasattr(shape, "text"):
					text += shape.text + "\n"
	else:
		raise ValueError("Unsupported file type for extraction")
	return text.strip()


def chunk_text(raw_text: str) -> List[str]:
	splitter = RecursiveCharacterTextSplitter(
		chunk_size=1500,
		chunk_overlap=200,
		length_function=len,
	)
	return splitter.split_text(raw_text)


def get_embeddings_model():
	model_name = os.getenv("EMBEDDINGS_MODEL", "sentence-transformers/all-MiniLM-L6-v2")
	return HuggingFaceEmbeddings(model_name=model_name)


# Minimal in-memory vector search for small docs (Windows-friendly)
_embedder = None


def _get_embedder():
	global _embedder
	if _embedder is None:
		_embedder = get_embeddings_model()
	return _embedder


def build_or_load_vectorstore(file_id: str, chunks: List[str]):
	# For development on Windows, we skip persistent vector DB
	# and keep data in memory for the request lifecycle.
	# Caller already has chunks; we just return them.
	return chunks


def retrieve_context(file_id: str, query: str, k: int = 6) -> List[Tuple[str, float]]:
	embedder = _get_embedder()
	# In-memory cosine similarity
	# Compute embeddings
	# Note: sentence-transformers returns lists of floats
	chunk_texts = build_or_load_vectorstore(file_id, chunks=[])
	if not chunk_texts:
		return []
	chunk_embs = embedder.embed_documents(chunk_texts)
	query_emb = embedder.embed_query(query)
	# Cosine similarity
	import numpy as np
	q = np.array(query_emb)
	scores = []
	for t, e in zip(chunk_texts, chunk_embs):
		v = np.array(e)
		den = (np.linalg.norm(q) * np.linalg.norm(v)) or 1.0
		score = float(np.dot(q, v) / den)
		scores.append((t, score))
	scores.sort(key=lambda x: x[1], reverse=True)
	return [(t, 1 - s) for t, s in scores[:k]]  # smaller is better distance


def _call_ollama(prompt: str, model: str = None) -> str:
	"""Call Ollama LLM via HTTP API or CLI fallback"""
	if model is None:
		model = os.getenv("OLLAMA_MODEL", "llama3.1:8b")
	
	# Try HTTP first (ollama serve at 127.0.0.1:11434)
	try:
		url = "http://127.0.0.1:11434/api/generate"
		payload = {
			"model": model, 
			"prompt": prompt, 
			"stream": False,
			"options": {
				"temperature": 0.7,
				"top_p": 0.9,
				"num_predict": 2048
			}
		}
		with httpx.Client(timeout=120) as client:
			resp = client.post(url, json=payload)
			if resp.status_code == 200:
				data = resp.json()
				return data.get("response", "")
			else:
				print(f"Ollama HTTP error: {resp.status_code} - {resp.text}")
	except Exception as e:
		print(f"Ollama HTTP request failed: {e}")
	
	# Fallback CLI
	try:
		result = subprocess.run(
			["ollama", "run", model, prompt], 
			capture_output=True, 
			text=True, 
			timeout=180
		)
		if result.returncode == 0:
			return result.stdout
		else:
			print(f"Ollama CLI error: {result.stderr}")
	except Exception as e:
		print(f"Ollama CLI failed: {e}")
	
	return ""


def summarize_document(chunks: List[str], max_length: int = 500) -> Dict[str, any]:
	"""Generate a summary and key points from document chunks using Ollama"""
	if not chunks:
		return {
			"summary": "No content to summarize.",
			"key_points": [],
			"word_count": 0
		}
	
	# Join chunks and limit context for LLM
	context = "\n\n".join(chunks[:10])  # Limit to first 10 chunks to avoid token limits
	
	prompt = f"""You are an expert document summarizer. Analyze the following document and provide:

1. A concise summary (maximum {max_length} words)
2. 5-7 key points or main ideas
3. The total word count of your summary

Format your response as valid JSON:
{{
  "summary": "Your summary text here...",
  "key_points": ["Key point 1", "Key point 2", "Key point 3"],
  "word_count": 123
}}

Document content:
{context}

Remember: Return ONLY valid JSON, no additional text before or after."""

	response_text = _call_ollama(prompt)
	
	if response_text:
		try:
			# Clean the response to extract JSON
			response_text = response_text.strip()
			if response_text.startswith("```json"):
				response_text = response_text[7:]
			if response_text.endswith("```"):
				response_text = response_text[:-3]
			
			payload = json.loads(response_text)
			if isinstance(payload, dict):
				return {
					"summary": payload.get("summary", "Summary generation failed."),
					"key_points": payload.get("key_points", []),
					"word_count": payload.get("word_count", 0)
				}
		except json.JSONDecodeError as e:
			print(f"Failed to parse Ollama JSON response: {e}")
			print(f"Raw response: {response_text[:200]}...")
	
	# Fallback if LLM fails
	fallback_summary = " ".join(chunks[0].split()[:max_length//5]) if chunks else "No content available."
	return {
		"summary": fallback_summary,
		"key_points": ["Content analysis failed - please check Ollama connection"],
		"word_count": len(fallback_summary.split())
	}


def generate_quiz_from_chunks(chunks: List[str], num_questions: int = 5) -> List[dict]:
	"""Generate quiz questions using Ollama LLM"""
	if not chunks:
		return []
	
	# Limit context to avoid token limits
	context = "\n\n".join(chunks[:8]) if len(chunks) > 8 else "\n\n".join(chunks)
	
	prompt = f"""You are an expert educator creating multiple-choice questions from study materials.

Create {num_questions} high-quality multiple-choice questions based on the provided content.

Requirements:
- Questions should be directly answerable from the given content
- Each question must have exactly 4 options (A, B, C, D)
- Provide clear, unambiguous questions
- Include brief explanations for correct answers
- Avoid external knowledge not present in the content

Format your response as valid JSON:
{{
  "questions": [
    {{
      "question": "What is the main topic discussed in this document?",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "answer": "A",
      "explanation": "The document clearly states...",
      "source_chunks": ["Relevant text snippet from the content"]
    }}
  ]
}}


Content to analyze:
{context}

Return ONLY valid JSON, no additional text."""

	response_text = _call_ollama(prompt)
	
	if response_text:
		try:
			# Clean the response to extract JSON
			response_text = response_text.strip()
			if response_text.startswith("```json"):
				response_text = response_text[7:]
			if response_text.endswith("```"):
				response_text = response_text[:-3]
			
			payload = json.loads(response_text)
			if isinstance(payload, dict) and isinstance(payload.get("questions"), list):
				questions = []
				for q in payload["questions"][:num_questions]:
					if all(key in q for key in ["question", "options", "answer", "explanation"]):
						questions.append({
							"question": q["question"],
							"options": q["options"][:4],  # Ensure exactly 4 options
							"answer": q["answer"],
							"explanation": q["explanation"],
							"source_chunks": q.get("source_chunks", [])
						})
				if questions:
					return questions
		except json.JSONDecodeError as e:
			print(f"Failed to parse quiz JSON response: {e}")
			print(f"Raw response: {response_text[:200]}...")
	
	# Fallback if LLM fails
	print("LLM quiz generation failed, using fallback questions")
	fallback_questions = []
	for i in range(min(num_questions, len(chunks))):
		context = chunks[i][:300] if chunks[i] else "No content available"
		fallback_questions.append({
			"question": f"Question {i+1}: Based on the provided material, what is the main point?",
			"options": ["Option A", "Option B", "Option C", "Option D"],
			"answer": "A",
			"explanation": "This is a fallback question. Please ensure Ollama is running for AI-generated questions.",
			"source_chunks": [context]
		})
	return fallback_questions

def generate_questions_from_chunks(chunks, num_questions):
    """
    Generate only questions (no options or answers)
    based on provided text chunks.
    """
    from ollama import chat  # or your existing model import
    combined_text = "\n".join(chunks)

    prompt = f"""
    Generate {num_questions} clear and concise **questions only** 
    based on the following content.
    Do NOT include answers or multiple-choice options.
    Number each question sequentially.

    Content:
    {combined_text}
    """

    try:
        response = chat(model="llama3", messages=[{"role": "user", "content": prompt}])
        output = response['message']['content']

        # Split into individual questions
        questions = [q.strip() for q in output.split("\n") if q.strip()]
        return questions

    except Exception as e:
        print(f"Error in generate_questions_from_chunks: {e}")
        return []
